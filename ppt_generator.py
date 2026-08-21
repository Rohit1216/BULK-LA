"""
ppt_generator.py
-----------------
Reads a multi-sheet LinkedIn-activity Excel file (one sheet per executive,
e.g. LA1, LA2, ... LA10) and fills the "Sample.pptx" template - producing
ONE PowerPoint that contains every executive's slides back-to-back, each
section with its own LinkedIn URL / Connections header and its own
"(1/2)"-style page counter.

This file has NO Streamlit code in it on purpose - it can be imported
and tested on its own, as well as used by app.py.
"""

import copy
import re
from io import BytesIO

import openpyxl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# 1. Settings you may want to tweak
# ---------------------------------------------------------------------------

# Column headers expected in each sheet's table (row can be anywhere in the
# first few rows - it's auto-detected by looking for "PostType").
EXCEL_COL_POST_TYPE = "PostType"
EXCEL_COL_TIME_AGO = "TimeAgo"
EXCEL_COL_CONTENT = "Content"
EXCEL_COL_POST_LINK = "PostLink"
EXCEL_COL_SOURCE = "Source"      # the cell that holds the hyperlink
EXCEL_COL_HEADLINE = "Headline"

# How many rows above the header row to scan for the "Linkedin URL" /
# "Connection" labels (column A) and their values (column B).
MAX_HEADER_SCAN_ROWS = 15

# Color used for the "Source" hyperlink text (typical PowerPoint link blue)
HYPERLINK_COLOR = RGBColor(0x05, 0x66, 0xC2)

# Placeholder text in Sample.pptx that marks the LinkedIn URL / Connections
# boxes, so we can find them regardless of shape naming.
LINKEDIN_ID_PLACEHOLDER_TEXT = "xxxx/"
CONNECTIONS_PLACEHOLDER_TEXT = "500+\n connections"

# The name/country text under the headshot photo has no data source (no
# "executive name" column in the sheet), so that text is cleared - but the
# photo itself and its card background are kept on every slide.
BRAND_NAME_PLACEHOLDER_TEXT = "xxxxxx"
BRAND_COUNTRY_PLACEHOLDER_TEXT = "xxxx"

# Shapes making up the header bar (title + LinkedIn ID/Connections boxes)
# and the data table - stripped out on the section-divider slide that's
# inserted before each executive's run of slides.
HEADER_BAR_SHAPE_NAMES = {"Rectangle 8", "Group 10", "Group 13"}


# ---------------------------------------------------------------------------
# 2. Read the Excel workbook (many sheets, one per executive)
# ---------------------------------------------------------------------------
def _find_header_row(ws):
    """Scans the first few rows for the one containing the 'PostType'
    column header, wherever it happens to sit."""
    for r in range(1, min(MAX_HEADER_SCAN_ROWS, ws.max_row) + 1):
        values = [str(c.value).strip() if c.value is not None else "" for c in ws[r]]
        if EXCEL_COL_POST_TYPE in values:
            return r
    return None


def read_sheet(ws):
    """
    Reads one worksheet and returns:
        {
          "sheet_name":  "LA1",
          "linkedin_url": "https://www.linkedin.com/in/...",
          "connections":  "500+\n connections",
          "rows": [ {post_type, timeline, headline, content, source_url}, ... ]
        }
    Returns None if the sheet has no "PostType" header anywhere in its
    first rows (so stray/blank sheets are skipped instead of erroring out).
    """
    header_row = _find_header_row(ws)
    if header_row is None:
        return None

    # LinkedIn URL / Connections sit in column A (label) / column B (value)
    # somewhere above the header row.
    linkedin_url = ""
    connections = ""
    for r in range(1, header_row):
        label = ws.cell(row=r, column=1).value
        value = ws.cell(row=r, column=2).value
        if not label:
            continue
        label_l = str(label).lower()
        if "linked" in label_l:
            linkedin_url = str(value).strip() if value else ""
        elif "connection" in label_l:
            connections = str(value).strip() if value else ""

    headers = {}
    for cell in ws[header_row]:
        if cell.value:
            headers[str(cell.value).strip()] = cell.column

    required = [EXCEL_COL_POST_TYPE, EXCEL_COL_TIME_AGO, EXCEL_COL_CONTENT,
                EXCEL_COL_SOURCE, EXCEL_COL_HEADLINE]
    missing = [c for c in required if c not in headers]
    if missing:
        raise ValueError(
            f"Sheet '{ws.title}': missing expected column headers: {missing}. "
            f"Found headers: {list(headers.keys())}"
        )

    rows = []
    for r in range(header_row + 1, ws.max_row + 1):
        post_type = ws.cell(row=r, column=headers[EXCEL_COL_POST_TYPE]).value
        timeline = ws.cell(row=r, column=headers[EXCEL_COL_TIME_AGO]).value
        content = ws.cell(row=r, column=headers[EXCEL_COL_CONTENT]).value
        headline = ws.cell(row=r, column=headers[EXCEL_COL_HEADLINE]).value

        if not any([post_type, timeline, content, headline]):
            continue

        source_cell = ws.cell(row=r, column=headers[EXCEL_COL_SOURCE])
        source_url = source_cell.hyperlink.target if source_cell.hyperlink else None
        if not source_url and EXCEL_COL_POST_LINK in headers:
            source_url = ws.cell(row=r, column=headers[EXCEL_COL_POST_LINK]).value

        rows.append({
            "post_type": post_type or "",
            "timeline": timeline or "",
            "headline": headline or "",
            "content": content or "",
            "source_url": source_url or "",
        })

    return {
        "sheet_name": ws.title,
        "linkedin_url": linkedin_url,
        "connections": connections,
        "rows": rows,
    }


def read_workbook(file_like):
    """Reads every sheet in the workbook (in tab order) and returns a list
    of sheet dicts, skipping any sheet that has no usable data."""
    wb = openpyxl.load_workbook(file_like, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        data = read_sheet(ws)
        if data and data["rows"]:
            sheets.append(data)
    if not sheets:
        raise ValueError(
            "No usable LinkedIn-activity data found in any sheet of the workbook. "
            "Each sheet needs a 'PostType' column header somewhere in its first rows."
        )
    return sheets


# ---------------------------------------------------------------------------
# 3. Slide duplication / deletion helpers (python-pptx has neither built-in)
# ---------------------------------------------------------------------------
def duplicate_slide(prs, slide_index):
    """
    Duplicates the slide at slide_index (0-based) and appends the copy at
    the end of the presentation. Returns the new slide.
    Handles pictures / relationships so images are not lost.
    """
    source = prs.slides[slide_index]
    dest = prs.slides.add_slide(source.slide_layout)

    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    rid_map = {}
    for rel_id, rel in source.part.rels.items():
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            new_rel_id = dest.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rel_id = dest.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel_id] = new_rel_id

    for shp in source.shapes:
        new_el = copy.deepcopy(shp._element)
        for el in new_el.iter():
            for attr in el.attrib:
                if el.attrib[attr] in rid_map:
                    el.attrib[attr] = rid_map[el.attrib[attr]]
        dest.shapes._spTree.append(new_el)

    return dest


def delete_slide(prs, slide_index):
    """Removes a slide entirely from the deck (slide + its relationship)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    id_to_remove = slides[slide_index].get(qn('r:id'))
    prs.part.drop_rel(id_to_remove)
    xml_slides.remove(slides[slide_index])


# ---------------------------------------------------------------------------
# 4. Shape / table helpers
# ---------------------------------------------------------------------------
def find_table_shape(slide):
    for shp in slide.shapes:
        if shp.has_table:
            return shp
    raise ValueError("No table found on the template slide.")


def iter_all_shapes(shapes):
    """Yield every shape, descending into groups."""
    for shp in shapes:
        yield shp
        if shp.shape_type == 6:  # GROUP
            yield from iter_all_shapes(shp.shapes)


def find_shape_by_exact_text(slide, text):
    for shp in iter_all_shapes(slide.shapes):
        if shp.has_text_frame and shp.text_frame.text.strip() == text.strip():
            return shp
    return None


def find_title_shape(slide):
    for shp in slide.shapes:
        if shp.has_text_frame and shp.name.lower().startswith("title"):
            return shp
    return None


def remove_shapes_by_name(slide, names):
    for shp in list(slide.shapes):
        if shp.name in names:
            shp._element.getparent().remove(shp._element)


def clear_unused_branding(slide):
    """The name/country text under the headshot has no data source, so it's
    blanked out (rather than left showing 'xxxxxx' / 'xxxx'). The photo and
    its card background are left exactly as in the template."""
    for placeholder in (BRAND_NAME_PLACEHOLDER_TEXT, BRAND_COUNTRY_PLACEHOLDER_TEXT):
        shp = find_shape_by_exact_text(slide, placeholder)
        if shp is not None:
            set_cell_text_frame(shp.text_frame, "")


def set_cell_text_frame(text_frame, text):
    """Replace the text of any text_frame (table cell or plain shape) while
    keeping the first run's formatting."""
    p = text_frame.paragraphs[0]
    if p.runs:
        run = p.runs[0]
        run.text = str(text)
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
        run.text = str(text)
    for extra_p in text_frame.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)


def _linkedin_display_id(linkedin_url):
    """The 'LINKEDIN ID' box is sized for a short handle (matches the
    template's own label), not a full URL, so pull just the profile slug
    out of the URL for display - e.g.
    'https://www.linkedin.com/in/jane-doe-12345/?x=1' -> 'jane-doe-12345/'.
    Falls back to the raw value if it doesn't look like a linkedin.com URL.
    """
    match = re.search(r"linkedin\.com/in/([^/?#]+)", linkedin_url, re.IGNORECASE)
    if match:
        return match.group(1) + "/"
    return linkedin_url


def set_header(slide, linkedin_url, connections):
    """Fills the LinkedIn URL and Connections boxes at the top of the slide
    straight from the Excel sheet. The LinkedIn box shows the short profile
    slug (to fit the template's box) but stays clickable to the full URL."""
    if linkedin_url:
        shp = find_shape_by_exact_text(slide, LINKEDIN_ID_PLACEHOLDER_TEXT)
        if shp is not None:
            shp.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            set_cell_text_frame(shp.text_frame, _linkedin_display_id(linkedin_url))
            run = shp.text_frame.paragraphs[0].runs[0]
            run.hyperlink.address = linkedin_url

    if connections:
        shp = find_shape_by_exact_text(slide, CONNECTIONS_PLACEHOLDER_TEXT)
        if shp is not None:
            set_cell_text_frame(shp.text_frame, connections)


def delete_table_row(table, row_idx):
    tbl = table._tbl
    trs = tbl.findall(qn('a:tr'))
    tbl.remove(trs[row_idx])


def set_cell_text(cell, text):
    """Replace the text of a cell while keeping the first run's formatting."""
    tf = cell.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        run = p.runs[0]
        run.text = str(text)
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        run = p.add_run()
        run.text = str(text)
    for extra_p in tf.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)


def set_details_cell(cell, content_text, source_url):
    """
    Fills the 'Details' cell with the post content, then appends a
    hyperlinked word 'Source' at the end (only if a URL is available).
    """
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(content_text)
    for extra in p.runs[1:]:
        extra._r.getparent().remove(extra._r)
    for extra_p in tf.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)

    if source_url:
        sep = p.add_run()
        sep.text = "  "
        sep.font.size = run.font.size
        sep.font.name = run.font.name

        link_run = p.add_run()
        link_run.text = "Source"
        link_run.font.size = run.font.size
        link_run.font.name = run.font.name
        link_run.font.underline = True
        link_run.font.color.rgb = HYPERLINK_COLOR
        link_run.hyperlink.address = source_url


# ---------------------------------------------------------------------------
# 5. Main entry point
# ---------------------------------------------------------------------------
def build_divider_slide(prs, sheet):
    """Inserts a section-break slide before an executive's run of content
    slides, so it's obvious in the deck where one executive ends and the
    next begins. Keeps the title banner and the photo card (blank name/
    country, since there's no data for them); strips the table and the
    LinkedIn ID / Connections header bar, since those belong to the
    content slides that follow."""
    slide = duplicate_slide(prs, 0)

    table_shape = find_table_shape(slide)
    table_shape._element.getparent().remove(table_shape._element)
    remove_shapes_by_name(slide, HEADER_BAR_SHAPE_NAMES)

    # Blank the country line (no data source for it). The name line becomes
    # a light identity tag for this section - the profile handle - instead
    # of staying blank, so the divider clearly says whose posts follow.
    country_shp = find_shape_by_exact_text(slide, BRAND_COUNTRY_PLACEHOLDER_TEXT)
    if country_shp is not None:
        set_cell_text_frame(country_shp.text_frame, "")

    name_shp = find_shape_by_exact_text(slide, BRAND_NAME_PLACEHOLDER_TEXT)
    if name_shp is not None:
        handle = _linkedin_display_id(sheet["linkedin_url"]) if sheet["linkedin_url"] else ""
        set_cell_text_frame(name_shp.text_frame, handle)

    title_shape = find_title_shape(slide)
    if title_shape is not None:
        for p in title_shape.text_frame.paragraphs:
            for r in p.runs:
                if re.search(r"\(\d+/\d+\)", r.text):
                    r.text = re.sub(r"\(\d+/\d+\)", f"— {sheet['sheet_name']}", r.text)

    return slide


def generate_pptx(template_path, excel_file_like, rows_per_slide=None):
    """
    template_path   : path to Sample.pptx
    excel_file_like : path or file-like object for the .xlsx/.xlsm upload.
                       Every sheet in the workbook that has a "PostType"
                       column is treated as one executive's data and gets
                       its own run of slides (own LinkedIn URL, Connections
                       and its own "(1/2)" page counter), appended one
                       after another into a single output deck.
    rows_per_slide  : how many table rows to put on each slide. Defaults to
                       however many data rows the template slide has (3).
    Returns: BytesIO of the finished .pptx
    """
    sheets = read_workbook(excel_file_like)

    prs = Presentation(template_path)
    # Slide 0 is kept as a pristine template - we only ever duplicate FROM
    # it (never fill it directly) so every sheet gets a fresh copy of the
    # untouched placeholders, then it is deleted at the very end.
    template_table_shape = find_table_shape(prs.slides[0])
    n_body_rows = len(template_table_shape.table.rows) - 1  # minus header row
    if rows_per_slide is None:
        rows_per_slide = n_body_rows

    def fill_new_slide(chunk, page_no, total_slides, linkedin_url, connections):
        slide = duplicate_slide(prs, 0)
        clear_unused_branding(slide)
        set_header(slide, linkedin_url, connections)

        table_shape = find_table_shape(slide)
        table = table_shape.table

        current_body_rows = len(table.rows) - 1
        while current_body_rows > len(chunk):
            delete_table_row(table, current_body_rows)
            current_body_rows -= 1
        # NOTE: if a chunk ever needs MORE rows than the template has,
        # increase rows_per_slide, or extend this to clone a row.

        for i, rec in enumerate(chunk):
            row_idx = i + 1  # +1 to skip header row
            set_cell_text(table.cell(row_idx, 0), rec["post_type"])
            set_cell_text(table.cell(row_idx, 1), rec["timeline"])
            set_cell_text(table.cell(row_idx, 2), rec["headline"])
            set_details_cell(table.cell(row_idx, 3), rec["content"], rec["source_url"])

        title_shape = find_title_shape(slide)
        if title_shape is not None:
            for p in title_shape.text_frame.paragraphs:
                for r in p.runs:
                    if re.search(r"\(\d+/\d+\)", r.text):
                        r.text = re.sub(r"\(\d+/\d+\)", f"({page_no}/{total_slides})", r.text)

        return slide

    for sheet in sheets:
        build_divider_slide(prs, sheet)
        records = sheet["rows"]
        chunks = [records[i:i + rows_per_slide] for i in range(0, len(records), rows_per_slide)]
        total_slides = len(chunks)
        for page_no, chunk in enumerate(chunks, start=1):
            fill_new_slide(chunk, page_no, total_slides, sheet["linkedin_url"], sheet["connections"])

    # The pristine template slide was never filled with real data - drop it.
    delete_slide(prs, 0)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out